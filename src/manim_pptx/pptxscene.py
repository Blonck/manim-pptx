from manim import *
import os

import manim
import pptx
import subprocess
import lxml.etree as etree
from functools import reduce

import logging

logger = logging.getLogger(__name__)
logger.setLevel(logging.DEBUG)


class PPTXScene(Scene):
    def __init__(self, *args, **kwargs):
        self.output_folder = kwargs.pop("output_folder", "./pptx/")
        self.temporary_dir = kwargs.pop("temporary_dir", "./temp/")
        super(PPTXScene, self).__init__(*args, **kwargs)

        self.slides = list()

        self.currentSlide = 1
        self.currentAnimation = 0
        self.slideStartAnimation = 0

    def construct(self, *args, **kwargs):
        return super().construct(*args, **kwargs)

    def play(self, *args, **kwargs):
        super(PPTXScene, self).play(*args, **kwargs)
        self.currentAnimation += 1
        logger.info(f"Add animation: {self.currentAnimation}")

    def wait(self, *args, **kwargs):
        super(PPTXScene, self).wait(*args, **kwargs)
        self.currentAnimation += 1

    def endSlide(self, loop=False, autonext=False, notes=None, shownextnotes=False):
        logger.info(f"End slide: {self.currentSlide} with animations [{self.slideStartAnimation},  {self.currentAnimation}]")
        self.slides.append(dict(
            type="loop" if loop else "slide",
            start=self.slideStartAnimation,
            end=self.currentAnimation,
            number=self.currentSlide,
            autonext=autonext,
            notes=notes,
            shownextnotes=shownextnotes,
        ))
        self.currentSlide += 1
        self.slideStartAnimation = self.currentAnimation

    def save_video_thumb(self, filename, imgname):
        subprocess.run([
            constants.FFMPEG_BIN,
            '-i', filename,
            '-vframes', '1',  # one frame
            '-loglevel', 'error',
            '-y',  # overwrite
            imgname,
        ], stdout=subprocess.PIPE)
        return imgname

    def get_dur(self, filename):
        return int(float(subprocess.check_output([
            "ffprobe",
            '-i', filename,
            "-show_entries", "format=duration", # show duration
            "-v", "quiet", # hide other outputs
            "-of", "csv=p=0", # only number
        ]).decode("utf-8").strip()) * 1000)

    def render(self, *args, **kwargs):
        super(PPTXScene, self).render(*args, **kwargs)

        if not os.path.exists(self.output_folder):
            os.mkdir(self.output_folder)

        if not os.path.exists(self.temporary_dir):
            os.mkdir(self.temporary_dir)

        logger.info("Creating PPTX")

        prs = pptx.Presentation(pptx=os.path.join(os.path.split(__file__)[0], "template.pptx"))

        prs.slide_width = self.camera.pixel_width * 9525 # pixels to emu
        prs.slide_height = self.camera.pixel_height * 9525

        blank_slide_layout = prs.slide_layouts[6]

        # for src_file in self.renderer.file_writer.partial_movie_files:
        #     print(src_file)

        #     thumb_file = os.path.join(self.temporary_dir, os.path.basename(src_file) + ".png")
        #     self.save_video_thumb(src_file, thumb_file)

        #     slide = prs.slides.add_slide(blank_slide_layout)

        #     # Add the video to the slide
        #     clip = slide.shapes.add_movie(src_file, 0, 0, prs.slide_width, prs.slide_height, mime_type='video/mp4', poster_frame_image=thumb_file)

        #     clipid = clip.element[0][0].attrib.get("id")

        #     # slide.shapes.add_movie(src_file, 0, 0, prs.slide_width, prs.slide_height)

        url_schema = "{http://schemas.openxmlformats.org/presentationml/2006/main}"

        for tslidei, tslide in enumerate(self.slides):
            logger.debug("Add slide {tslidei} with animations [{tslide['start']}, tslide['end']]")

            slide_movie_files = self.renderer.file_writer.partial_movie_files[tslide["start"]:tslide["end"]]

            slide = prs.slides.add_slide(blank_slide_layout)

            notes = tslide["notes"] if tslide["notes"] else ""

            if tslide["shownextnotes"] and len(self.slides) > tslidei + 1:
                notes += "\n" + "\n".join(list(map(lambda x: "> " + x, self.slides[tslidei + 1]["notes"].split("\n"))))

            slide.notes_slide.notes_text_frame.text = notes

            pics = list()

            for src_file in slide_movie_files:
                thumb_file = os.path.join(self.temporary_dir, os.path.basename(src_file) + ".png")
                self.save_video_thumb(src_file, thumb_file)

                logger.debug(f"adding video {src_file}")
                clip = slide.shapes.add_movie(src_file, 0, 0, prs.slide_width, prs.slide_height, mime_type='video/mp4', poster_frame_image=thumb_file)


                pics.append(dict(
                    id=clip.element[0][0].attrib.get("id"),
                    dur=self.get_dur(src_file),
                ))

            def addAutoNext():
                transition = etree.Element(url_schema + "transition", {
                    "spd":"slow",
                    "advTm":"0",
                })
                slide.element.insert(2, transition)

            if len(pics) > 0:

                outerchildTnLst = slide.element[2][0][0][0][0]

                if tslide["autonext"]:
                    addAutoNext()
                    outerchildTnLst = slide.element[3][0][0][0][0]

                seq = etree.Element(url_schema + "seq", concurrent="1", nextAc="seek")
                outerchildTnLst.insert(0,seq)


                def getcTnIDCounter():
                    getcTnIDCounter.cTnIDCounter += 1
                    return getcTnIDCounter.cTnIDCounter

                getcTnIDCounter.cTnIDCounter = 1

                def addCTn():
                    innercTn = etree.Element(url_schema + "cTn", id=str(getcTnIDCounter()), dur="indefinite", nodeType="mainSeq")
                    childTnLst = etree.Element(url_schema + "childTnLst")
                    par1 = etree.Element(url_schema + "par")
                    cTn = etree.Element(url_schema + "cTn", id=str(getcTnIDCounter()), fill="hold")
                    if tslide["type"] == "loop":
                        cTn.attrib["dur"] = str(reduce(lambda x, y: x + y, [p["dur"] for p in pics]))
                        cTn.attrib["repeatCount"] = "indefinite"
                    stCondLst = etree.Element(url_schema + "stCondLst")
                    cond1 = etree.Element(url_schema + "cond", delay="indefinite")
                    cond2 = etree.Element(url_schema + "cond", evt="onBegin", delay="0")
                    cond2tn = etree.Element(url_schema + "tn", val="2")
                    cond2.append(cond2tn)
                    stCondLst.append(cond1)
                    stCondLst.append(cond2)
                    cTn.append(stCondLst)
                    childTnLst2 = etree.Element(url_schema + "childTnLst")
                    cTn.append(childTnLst2)

                    par1.append(cTn)
                    childTnLst.append(par1)
                    innercTn.append(childTnLst)
                    seq.append(innercTn)
                    return childTnLst2

                childTnLst = addCTn()

                def addPrevCondLst():
                    prevCondLst = etree.Element(url_schema + "prevCondLst")
                    cond = etree.Element(url_schema + "cond", evt="onPrev", delay="0")
                    tgtEl = etree.Element(url_schema + "tgtEl")
                    sldTgt = etree.Element(url_schema + "sldTgt")
                    tgtEl.append(sldTgt)
                    cond.append(tgtEl)
                    prevCondLst.append(cond)
                    seq.append(prevCondLst)

                def addNextCondLst():
                    nextCondLst = etree.Element(url_schema + "nextCondLst")
                    cond = etree.Element(url_schema + "cond", evt="onNext", delay="0")
                    tgtEl = etree.Element(url_schema + "tgtEl")
                    sldTgt = etree.Element(url_schema + "sldTgt")
                    tgtEl.append(sldTgt)
                    cond.append(tgtEl)
                    nextCondLst.append(cond)
                    seq.append(nextCondLst)

                addPrevCondLst()
                addNextCondLst()

                currentdelay = 0
                for i, pic in enumerate(pics):
                    def addToFrontEffect():
                        par = etree.Element(url_schema + "par")
                        cTn = etree.Element(url_schema + "cTn", id=str(getcTnIDCounter()), fill="hold")
                        stCondLst = etree.Element(url_schema + "stCondLst")
                        cond = etree.Element(url_schema + "cond", delay=str(currentdelay))
                        stCondLst.append(cond)
                        cTn.append(stCondLst)

                        innerchildTnLst = etree.Element(url_schema + "childTnLst")
                        innerPar = etree.Element(url_schema + "par")
                        innercTn = etree.Element(url_schema + "cTn", id=str(getcTnIDCounter()), presetID="1", presetClass="entr", presetSubtype="0", fill="hold", nodeType="withEffect" if i == 0 else "afterEffect")
                        innerstCondLst = etree.Element(url_schema + "stCondLst")
                        innercond = etree.Element(url_schema + "cond", delay="0")
                        innerstCondLst.append(innercond)
                        innercTn.append(innerstCondLst)

                        innerInnerChildTnLst = etree.Element(url_schema + "childTnLst")
                        innercTn.append(innerInnerChildTnLst)

                        set = etree.Element(url_schema + "set")
                        cBhvr = etree.Element(url_schema + "cBhvr")

                        cBhvrcTn = etree.Element(url_schema + "cTn", id=str(getcTnIDCounter()), dur="1", fill="hold")
                        cBhvrcTnstCondLst = etree.Element(url_schema + "stCondLst")
                        cBhvrcTncond = etree.Element(url_schema + "cond", delay="0")
                        cBhvrcTnstCondLst.append(cBhvrcTncond)
                        cBhvrcTn.append(cBhvrcTnstCondLst)
                        cBhvr.append(cBhvrcTn)
                        tgtEl = etree.Element(url_schema + "tgtEl")
                        spTgt = etree.Element(url_schema + "spTgt", spid=str(pic["id"]))
                        tgtEl.append(spTgt)
                        cBhvr.append(tgtEl)
                        attrNameLst = etree.Element(url_schema + "attrNameLst")
                        attrName = etree.Element(url_schema + "attrName")
                        attrName.text = "style.visibility"
                        attrNameLst.append(attrName)
                        cBhvr.append(attrNameLst)

                        set.append(cBhvr)
                        to = etree.Element(url_schema + "to")
                        strVal = etree.Element(url_schema + "strVal", val="visible")
                        to.append(strVal)
                        set.append(to)

                        innerInnerChildTnLst.append(set)


                        innerPar.append(innercTn)
                        innerchildTnLst.append(innerPar)
                        cTn.append(innerchildTnLst)

                        par.append(cTn)
                        childTnLst.append(par)

                    def addToBackEffect():
                        par = etree.Element(url_schema + "par")
                        cTn = etree.Element(url_schema + "cTn", id=str(getcTnIDCounter()), fill="hold")
                        stCondLst = etree.Element(url_schema + "stCondLst")
                        cond = etree.Element(url_schema + "cond", delay=str(currentdelay))
                        stCondLst.append(cond)
                        cTn.append(stCondLst)

                        innerchildTnLst = etree.Element(url_schema + "childTnLst")
                        innerPar = etree.Element(url_schema + "par")
                        innercTn = etree.Element(url_schema + "cTn", id=str(getcTnIDCounter()), presetID="1", presetClass="exit", presetSubtype="0", fill="hold", nodeType="afterEffect")
                        innerstCondLst = etree.Element(url_schema + "stCondLst")
                        innercond = etree.Element(url_schema + "cond", delay="0")
                        innerstCondLst.append(innercond)
                        innercTn.append(innerstCondLst)

                        innerInnerChildTnLst = etree.Element(url_schema + "childTnLst")
                        innercTn.append(innerInnerChildTnLst)

                        set = etree.Element(url_schema + "set")
                        cBhvr = etree.Element(url_schema + "cBhvr")

                        cBhvrcTn = etree.Element(url_schema + "cTn", id=str(getcTnIDCounter()), dur="1", fill="hold")
                        cBhvrcTnstCondLst = etree.Element(url_schema + "stCondLst")
                        cBhvrcTncond = etree.Element(url_schema + "cond", delay="0")
                        cBhvrcTnstCondLst.append(cBhvrcTncond)
                        cBhvrcTn.append(cBhvrcTnstCondLst)
                        cBhvr.append(cBhvrcTn)
                        tgtEl = etree.Element(url_schema + "tgtEl")
                        spTgt = etree.Element(url_schema + "spTgt", spid=str(pic["id"]))
                        tgtEl.append(spTgt)
                        cBhvr.append(tgtEl)
                        attrNameLst = etree.Element(url_schema + "attrNameLst")
                        attrName = etree.Element(url_schema + "attrName")
                        attrName.text = "style.visibility"
                        attrNameLst.append(attrName)
                        cBhvr.append(attrNameLst)

                        set.append(cBhvr)
                        to = etree.Element(url_schema + "to")
                        strVal = etree.Element(url_schema + "strVal", val="hidden")
                        to.append(strVal)
                        set.append(to)

                        innerInnerChildTnLst.append(set)

                        cmd = etree.Element(url_schema + "cmd", type="call", cmd="stop")
                        cBhvr = etree.Element(url_schema + "cBhvr")
                        cBhvrcTn = etree.Element(url_schema + "cTn", id=str(getcTnIDCounter()), dur="1", fill="hold")
                        cBhvrcTnstCondLst = etree.Element(url_schema + "stCondLst")
                        cBhvrcTnstCondLstCond = etree.Element(url_schema + "cond", delay="0")
                        cBhvrcTnstCondLst.append(cBhvrcTnstCondLstCond)
                        cBhvrcTn.append(cBhvrcTnstCondLst)
                        cBhvr.append(cBhvrcTn)
                        tgtEl = etree.Element(url_schema + "tgtEl")
                        spTgt = etree.Element(url_schema + "spTgt", spid=str(pic["id"]))
                        tgtEl.append(spTgt)
                        cBhvr.append(tgtEl)
                        cmd.append(cBhvr)

                        innerInnerChildTnLst.append(cmd)


                        innerPar.append(innercTn)
                        innerchildTnLst.append(innerPar)
                        cTn.append(innerchildTnLst)

                        par.append(cTn)
                        childTnLst.append(par)

                    def playEffect():
                        par = etree.Element(url_schema + "par")
                        cTn = etree.Element(url_schema + "cTn", id=str(getcTnIDCounter()), fill="hold")
                        stCondLst = etree.Element(url_schema + "stCondLst")
                        cond = etree.Element(url_schema + "cond", delay=str(currentdelay))
                        stCondLst.append(cond)
                        cTn.append(stCondLst)

                        innerchildTnLst = etree.Element(url_schema + "childTnLst")
                        innerPar = etree.Element(url_schema + "par")
                        innercTn = etree.Element(url_schema + "cTn", id=str(getcTnIDCounter()), presetID="1", presetClass="mediacall", presetSubtype="0", fill="hold", nodeType="afterEffect")
                        innerstCondLst = etree.Element(url_schema + "stCondLst")
                        innercond = etree.Element(url_schema + "cond", delay="0")
                        innerstCondLst.append(innercond)
                        innercTn.append(innerstCondLst)

                        innerInnerChildTnLst = etree.Element(url_schema + "childTnLst")
                        innercTn.append(innerInnerChildTnLst)

                        cmd = etree.Element(url_schema + "cmd", type="call", cmd="playFrom(0.0)")
                        cBhvr = etree.Element(url_schema + "cBhvr")
                        cBhvrcTn = etree.Element(url_schema + "cTn", id=str(getcTnIDCounter()), dur=str(pic["dur"]), fill="hold")
                        cBhvr.append(cBhvrcTn)
                        tgtEl = etree.Element(url_schema + "tgtEl")
                        spTgt = etree.Element(url_schema + "spTgt", spid=str(pic["id"]))
                        tgtEl.append(spTgt)
                        cBhvr.append(tgtEl)
                        cmd.append(cBhvr)

                        innerInnerChildTnLst.append(cmd)

                        innerPar.append(innercTn)
                        innerchildTnLst.append(innerPar)
                        cTn.append(innerchildTnLst)

                        par.append(cTn)
                        childTnLst.append(par)

                    addToFrontEffect()
                    playEffect()
                    currentdelay+=pic["dur"]
                    if i+1 != len(pics):# or tslide["type"] == "loop":
                        addToBackEffect()


                for i in range(1, len(outerchildTnLst)):
                    outerchildTnLst[i][0][0].attrib["id"] = str(getcTnIDCounter())

            # if len(outerchildTnLst) > 1:
            #     seq[0][0][0][0][0][1][0].attrib["val"] = outerchildTnLst[1][0][0].attrib["id"]


        prs.save(os.path.join(self.output_folder, type(self).__name__ + '.pptx'))
